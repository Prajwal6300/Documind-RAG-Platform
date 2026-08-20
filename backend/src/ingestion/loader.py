"""Multi-format document parsing and layout extraction for DocuMind.

Supports structure-preserving extraction for:
- PDF (PyMuPDF with page numbers and table detection)
- Word DOCX (python-docx paragraphs and tables)
- Plain Text & Markdown (TXT, MD)
- CSV (Delimited rows and columns)
- Excel XLSX (openpyxl sheet and table extraction)
- PowerPoint PPTX (python-pptx slides and text boxes)
"""

import csv
from pathlib import Path
import pymupdf
from docx import Document


def _format_table(rows: list) -> str:
    """Format a table into a structure-preserving textual block."""
    if not rows:
        return ""

    lines = []
    for row in rows:
        cells = [str(cell).strip() if cell is not None else "" for cell in row]
        if not any(cells):
            continue
        lines.append(" | ".join(cells))

    return "\n".join(lines)


def load_pdf(file_path: str) -> list[dict]:
    """Extract pages and tables from PDF documents.

    Raises ValueError with a clear message for password-protected or corrupt
    files instead of letting a cryptic exception surface.
    """
    pages = []
    try:
        pdf = pymupdf.open(file_path)
    except pymupdf.PdfBaseError as exc:
        raise ValueError(
            "The PDF could not be opened. It may be corrupted or password-protected."
        ) from exc
    except Exception as exc:
        raise ValueError("The PDF could not be read.") from exc

    try:
        if pdf.needs_pass:
            raise ValueError("The PDF is password-protected. Please provide an unprotected copy.")
    except ValueError:
        raise
    except Exception:
        pass

    try:
        for page_number, page in enumerate(pdf):
            text = page.get_text()
            table_blocks = []

            try:
                found = page.find_tables()
            except Exception:
                found = None

            if found and getattr(found, "tables", None):
                for table in found.tables:
                    try:
                        rows = table.extract()
                        block = _format_table(rows)
                        if block and block not in table_blocks:
                            table_blocks.append(block)
                    except Exception:
                        continue

            if table_blocks:
                text += "\n\nTABLE DATA:\n" + "\n\n".join(table_blocks)

            if text.strip():
                pages.append({
                    "text": text,
                    "page": page_number + 1
                })
    finally:
        try:
            pdf.close()
        except Exception:
            pass

    return pages


def load_docx(file_path: str) -> list[dict]:
    """Extract paragraphs and tables from Word documents."""
    try:
        document = Document(file_path)
    except Exception as exc:
        raise ValueError(
            "The DOCX file could not be opened. It may be corrupted or password-protected."
        ) from exc

    text_parts = [
        paragraph.text
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]

    for table in document.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])

        block = _format_table(rows)
        if block:
            text_parts.append("\nTABLE DATA:\n" + block)

    text = "\n".join(text_parts)

    return [{
        "text": text,
        "page": None
    }]


def load_txt(file_path: str) -> list[dict]:
    """Load plain text and markdown documents."""
    text = Path(file_path).read_text(
        encoding="utf-8",
        errors="ignore"
    )
    return [{
        "text": text,
        "page": None
    }]


def load_csv(file_path: str) -> list[dict]:
    """Load CSV rows into structured table text."""
    rows = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(row)
    except Exception:
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        return [{"text": text, "page": None}]

    table_block = _format_table(rows)
    return [{
        "text": "CSV DATA:\n" + table_block,
        "page": None
    }]


def load_xlsx(file_path: str) -> list[dict]:
    """Extract sheets and tabular data from XLSX files."""
    pages = []
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
    except Exception as exc:
        raise ValueError(
            "The XLSX file could not be parsed. It may be corrupted, password-protected, or not a valid Excel file."
        ) from exc

    try:
        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(v is not None and str(v).strip() for v in row):
                    rows.append([str(v).strip() if v is not None else "" for v in row])
            if rows:
                table_block = _format_table(rows)
                pages.append({
                    "text": f"SHEET: {sheet_name}\n\nTABLE DATA:\n" + table_block,
                    "page": sheet_idx + 1,
                })
    finally:
        try:
            wb.close()
        except Exception:
            pass

    if not pages:
        return [{"text": "", "page": None}]
    return pages


def load_pptx(file_path: str) -> list[dict]:
    """Extract slides and text boxes / tables from PPTX presentations."""
    pages = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as exc:
        raise ValueError(
            "The PPTX file could not be parsed. It may be corrupted or not a valid PowerPoint file."
        ) from exc

    try:
        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            slide_texts.append(t)
                elif shape.has_table:
                    t_rows = []
                    for row in shape.table.rows:
                        t_rows.append([cell.text.strip() for cell in row.cells])
                    table_block = _format_table(t_rows)
                    if table_block:
                        slide_texts.append("\nTABLE DATA:\n" + table_block)
            if slide_texts:
                pages.append({
                    "text": f"SLIDE {slide_idx + 1}:\n" + "\n".join(slide_texts),
                    "page": slide_idx + 1,
                })
    except Exception:
        pass

    if not pages:
        return [{"text": "", "page": None}]
    return pages


def load_document(file_path: str) -> list[dict]:
    """Unified document loader routing by file extension."""
    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        return load_pdf(file_path)

    if extension == ".docx":
        return load_docx(file_path)

    if extension in (".txt", ".md", ".markdown"):
        return load_txt(file_path)

    if extension == ".csv":
        return load_csv(file_path)

    if extension in (".xlsx", ".xls"):
        return load_xlsx(file_path)

    if extension in (".pptx", ".ppt"):
        return load_pptx(file_path)

    raise ValueError(
        f"Unsupported file type: {extension}"
    )
